#!/usr/bin/env python3
"""Personalise and extend the Heimdall Executive Briefing (Nordic) deck.

This script is *idempotent*: it always reads the pristine backup
``Heimdall_Executive_Briefing_Nordic.ORIGINAL.pptx`` and writes the working
``Heimdall_Executive_Briefing_Nordic.pptx``. Re-running never compounds edits.

What it does
------------
1. Applies a personal "template" treatment to every content slide — an author
   monogram chip ("JRP"), name + brand footer and page number — so the deck
   reads as a personal, signed piece of work rather than an anonymous template.
2. Inserts eight new slides, built in the exact house style (Segoe UI, navy
   #0A2133 cards, teal/blue spectrum bar), at the right points in the narrative:
     A  Inputs & outputs — what goes in, what comes out, how it flows
     B  Where Heimdall fits — in front of the core banking / mainframe
     C  Threat coverage — detection, false positives, false negatives, accuracy, throughput
     D  Elastic scale — absorbing a transaction spike
     E  Resilience — what happens when a component fails
     F  Security, enforced — how each control is guaranteed
     G  Operations dashboard — the live management view
     H  Roadmap & milestones — a 12-month path with exit criteria

Requires: python-pptx (pip install python-pptx)
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DOCS = "/mnt/c/Users/jpontvianne/Documents/Azure/MasterArchitect/Documents"
SRC = os.path.join(DOCS, "Heimdall_Executive_Briefing_Nordic.ORIGINAL.pptx")
OUT = os.path.join(DOCS, "Heimdall_Executive_Briefing_Nordic.pptx")

# --- House palette (sampled from the original deck) ------------------------- #
NAVY = "0A2133"
CARD = "0E2C44"
CARD2 = "0B2236"
BORDER = "21506F"
DIV = "21506F"
TEAL = "3DD6C4"
BLUE = "1F8FE5"
PURPLE = "8C7BE0"
GREEN = "5FD39B"
AMBER = "F2B441"
RED = "E5687A"
HEAD = "8FC9E8"
TITLE = "EAF6FF"
BODY = "9FB9CD"
EYEBROW = "3DD6C4"
FONT = "Segoe UI"

MID = MSO_ANCHOR.MIDDLE
TOP = MSO_ANCHOR.TOP
LEFT = PP_ALIGN.LEFT
CENTER = PP_ALIGN.CENTER

AUTHOR = "Jean-Rémi Pontvianne"
ROLE = "Azure Master Architect"
BRAND = "Heimdall · Real-Time Fraud Intelligence"


def rgb(h: str) -> RGBColor:
    return RGBColor.from_string(h)


def _lerp(a, b, t):
    return int(round(a + (b - a) * t))


def spectrum_color(t: float) -> str:
    """Teal -> blue -> deep-blue gradient, matching the original top bar."""
    teal = (0x3D, 0xD6, 0xC4)
    blue = (0x1F, 0x8F, 0xE5)
    deep = (0x1F, 0x4E, 0x78)
    if t <= 0.5:
        u = t / 0.5
        c = tuple(_lerp(teal[i], blue[i], u) for i in range(3))
    else:
        u = (t - 0.5) / 0.5
        c = tuple(_lerp(blue[i], deep[i], u) for i in range(3))
    return "%02X%02X%02X" % c


# --------------------------------------------------------------------------- #
# Low-level shape helpers
# --------------------------------------------------------------------------- #
def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(lw)
    return sp


def text(s, l, t, w, h, lines, size=13, color=BODY, bold=False, align=LEFT,
         anchor=TOP, space=4, font=FONT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        ov = {}
        if isinstance(ln, tuple):
            ln, ov = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.space_after = Pt(ov.get("space", space))
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name = font
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = rgb(ov.get("color", color))
    return tb


def spectrum(s):
    for i in range(61):
        rect(s, i * 0.2167, 0.0, 0.24, 0.13, fill=spectrum_color(i / 60.0))


def base_slide(prs, eyebrow, title, bg=NAVY):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    rect(s, 0, 0, 13.33, 7.5, fill=bg)
    spectrum(s)
    text(s, 0.7, 0.5, 11.9, 0.4, eyebrow, size=13, bold=True, color=EYEBROW)
    text(s, 0.7, 0.82, 11.9, 0.9, title, size=30, bold=True, color=HEAD)
    rect(s, 0.7, 1.62, 11.93, 0.03, fill=DIV)
    return s


def column_card(s, l, t, w, h, accent, heading, bullets, body_size=13):
    rect(s, l, t, w, h, fill=CARD, line=BORDER, lw=1.2, rounded=True)
    rect(s, l + 0.12, t, w - 0.24, 0.06, fill=accent)
    text(s, l + 0.32, t + 0.28, w - 0.6, 0.6, heading, size=16.5, bold=True, color=TITLE)
    rect(s, l + 0.32, t + 0.9, 0.55, 0.04, fill=accent)
    text(s, l + 0.32, t + 1.08, w - 0.6, h - 1.25,
         [("•  " + b, {}) for b in bullets], size=body_size, color=BODY, space=6)


def kpi_card(s, l, t, w, accent, value, label, desc, h=2.0):
    rect(s, l, t, w, h, fill=CARD, line=BORDER, lw=1.2, rounded=True)
    rect(s, l + 0.12, t, w - 0.24, 0.06, fill=accent)
    text(s, l + 0.05, t + 0.18, w - 0.1, 0.7, value, size=26, bold=True,
         color=accent, align=CENTER)
    text(s, l + 0.1, t + 0.92, w - 0.2, 0.4, label, size=13.5, bold=True,
         color=TITLE, align=CENTER)
    text(s, l + 0.12, t + 1.3, w - 0.24, h - 1.35, desc, size=10.5, color=BODY,
         align=CENTER)


def table_block(s, x, y, w, cols, rows, accents=None, row_h=0.66, head_h=0.4):
    """cols: list of (label, fraction). rows: list of lists of strings."""
    widths = [w * f for _, f in cols]
    cy = y
    rect(s, x, cy, w, head_h, fill=BORDER)
    cx = x
    for (label, _), cw in zip(cols, widths):
        text(s, cx + 0.14, cy, cw - 0.24, head_h, label, size=11, bold=True,
             color=TITLE, anchor=MID)
        cx += cw
    cy += head_h
    for ri, row in enumerate(rows):
        fill = CARD if ri % 2 == 0 else CARD2
        rect(s, x, cy, w, row_h, fill=fill, line=BORDER, lw=0.6)
        if accents:
            rect(s, x, cy, 0.06, row_h, fill=accents[ri])
        cx = x
        for ci, (cell, cw) in enumerate(zip(row, widths)):
            text(s, cx + 0.16, cy, cw - 0.28, row_h, cell,
                 size=11.5, color=(TITLE if ci == 0 else BODY),
                 bold=(ci == 0), anchor=MID, space=2)
            cx += cw
        cy += row_h
    return cy


def chip(s, l, t, w, label, fill, fg=NAVY, size=11):
    rect(s, l, t, w, 0.34, fill=fill, rounded=True)
    text(s, l, t + 0.02, w, 0.3, label, size=size, bold=True, color=fg,
         align=CENTER, anchor=MID)


def arrow(s, l, t, color=TEAL):
    text(s, l, t, 0.4, 0.4, "›", size=26, bold=True, color=color, align=CENTER,
         anchor=MID)


# --------------------------------------------------------------------------- #
# New slide builders
# --------------------------------------------------------------------------- #
def slide_inputs_outputs(prs):
    s = base_slide(prs, "HOW INFORMATION FLOWS", "What goes in, what comes out")
    # INPUTS card
    column_card(s, 0.7, 2.0, 4.0, 3.5, BLUE, "Inputs — what we receive", [
        "The payment request: card token, amount, currency, merchant, country, channel, device",
        "Context we already hold: the card's recent behaviour and the merchant's risk",
        "The network around them: who this card and merchant connect to (the graph)",
    ], body_size=12.5)
    text(s, 0.7, 5.55, 4.0, 0.5, "Plain English: \"Here is a payment — is it safe?\"",
         size=11.5, color=BODY)
    # ENGINE
    rect(s, 5.05, 2.55, 3.23, 2.4, fill="0E3350", line=TEAL, lw=1.4, rounded=True)
    text(s, 5.05, 2.75, 3.23, 0.5, "HEIMDALL", size=15, bold=True, color=TEAL,
         align=CENTER)
    text(s, 5.15, 3.25, 3.03, 1.6, [
        ("Enrich → Score (AI) → Decide", {"bold": True, "color": TITLE, "size": 12.5, "align": CENTER, "space": 8}),
        ("Ensemble + Graph Neural Net", {"align": CENTER, "size": 11, "space": 3}),
        ("~14 ms, in-line, explainable", {"align": CENTER, "size": 11, "color": TEAL}),
    ], align=CENTER)
    arrow(s, 4.6, 3.55)
    arrow(s, 8.28, 3.55)
    # OUTPUTS card
    column_card(s, 8.63, 2.0, 4.0, 3.5, TEAL, "Outputs — what we return", [
        "A decision in ~14 ms: Approve · Step-up (SCA) · Decline",
        "A plain reason + machine reason-codes for every call",
        "A risk score 0–1 and the PSD2 exemption applied",
        "An event for dashboards, learning and the EBA report",
    ], body_size=12.5)
    text(s, 8.63, 5.55, 4.0, 0.5, "Plain English: \"Yes / step up / no — and why.\"",
         size=11.5, color=BODY)
    # Flow strip
    steps = ["Transaction in", "Enrich features", "Score (AI)", "Decide & explain", "Emit & learn"]
    fw = 2.3
    x = 0.7
    y = 6.15
    for i, st in enumerate(steps):
        rect(s, x, y, fw, 0.5, fill=CARD, line=BORDER, lw=0.8, rounded=True)
        text(s, x, y + 0.02, fw, 0.46, f"{i+1}  {st}", size=11, bold=True,
             color=TITLE, align=CENTER, anchor=MID)
        if i < len(steps) - 1:
            arrow(s, x + fw - 0.02, y + 0.05)
        x += fw + 0.18
    return s


def slide_where_fits(prs):
    s = base_slide(prs, "WHERE IT FITS",
                   "In line, in front of your core — not in your way")
    stack = [
        ("Cardholder\n& merchant", BORDER, False),
        ("Acquirer /\ngateway", BORDER, False),
        ("HEIMDALL\nin-line scoring", TEAL, True),
        ("Scheme /\nissuer auth", BORDER, False),
        ("Core banking /\nmainframe", BLUE, False),
    ]
    bw = 2.18
    x = 0.7
    y = 2.25
    for i, (lbl, accent, hot) in enumerate(stack):
        fill = "0E3350" if hot else CARD
        sp = rect(s, x, y, bw, 1.5, fill=fill, line=accent, lw=1.6 if hot else 1.0,
                  rounded=True)
        rect(s, x + 0.1, y, bw - 0.2, 0.06, fill=accent)
        text(s, x + 0.06, y, bw - 0.12, 1.5, lbl.split("\n"),
             size=12.5, bold=hot, color=(TEAL if hot else TITLE), align=CENTER,
             anchor=MID, space=2)
        if i < len(stack) - 1:
            arrow(s, x + bw - 0.04, y + 0.55, color=TEAL)
        x += bw + 0.2
    text(s, 0.7, 3.95, 11.9, 0.5,
         "Heimdall sits as a decision point just before authorisation. The mainframe stays the system of record.",
         size=13, color=HEAD, bold=True)
    # three notes
    notes = [
        (TEAL, "A call in the auth path",
         ["One API call, a few ms added", "No rip-and-replace of the core", "Issuer, acquirer or processor"]),
        (BLUE, "Reads, never blocks",
         ["Async events feed the learning path", "Mainframe is never on the hot path", "Fails safe if a backend is down"]),
        (PURPLE, "Three ways to deploy",
         ["In-line — synchronous gate", "Sidecar — shadow / advisory", "Out-of-band — post-auth case work"]),
    ]
    cw = 3.85
    x = 0.7
    for accent, h, items in notes:
        column_card(s, x, 4.6, cw, 2.15, accent, h, items, body_size=12)
        x += cw + 0.19
    return s


def slide_threats(prs):
    s = base_slide(prs, "THREAT COVERAGE",
                   "Every known threat — and what we do with each call")
    # threat chips row
    threats = ["Card-not-present", "Account takeover", "Synthetic identity",
               "Fan-out / mule rings", "Card testing", "Merchant collusion",
               "Cross-border laundering"]
    x = 0.7
    y = 1.85
    for th in threats:
        w = 0.18 + len(th) * 0.085
        chip(s, x, y, w, th, TEAL, fg=NAVY, size=10.5)
        x += w + 0.14
    # 2x2 quadrant
    quad = [
        (GREEN, "True positive — fraud caught", "0.7", "2.45",
         ["Decline or step-up in-line", "Open a case, feed the graph"]),
        (AMBER, "False positive — good flagged", "6.72", "2.45",
         ["Graph context + PSD2 optimiser cut it", "Step-up, not hard-decline; HITL + retrain"]),
        (BLUE, "True negative — good approved", "0.7", "4.0",
         ["Frictionless, ~14 ms", "No challenge, no friction"]),
        (RED, "False negative — fraud missed", "6.72", "4.0",
         ["Graph + agents catch it within seconds", "Reflector + drift watch trigger retrain"]),
    ]
    for accent, h, qx, qy, items in quad:
        column_card(s, float(qx), float(qy), 5.9, 1.45, accent, h, items, body_size=11.5)
    # KPI strip
    kpis = [(GREEN, "0.987", "Model AUC", "back-tested"),
            (BLUE, "94%", "Precision", "of flags are fraud"),
            (TEAL, "91%", "Recall", "of fraud caught"),
            (AMBER, "1.1%", "False-positive", "down from 2.8%"),
            (PURPLE, "5k–20k", "TPS", "sustained → peak"),
            (HEAD, "<18 ms", "p99 latency", "in-line decision")]
    x = 0.7
    cw = 1.92
    for accent, v, l, d in kpis:
        rect(s, x, 5.62, cw, 1.1, fill=CARD, line=BORDER, lw=1.0, rounded=True)
        rect(s, x + 0.1, 5.62, cw - 0.2, 0.05, fill=accent)
        text(s, x, 5.74, cw, 0.45, v, size=18, bold=True, color=accent, align=CENTER)
        text(s, x, 6.18, cw, 0.3, l, size=10.5, bold=True, color=TITLE, align=CENTER)
        text(s, x, 6.45, cw, 0.25, d, size=9, color=BODY, align=CENTER)
        x += cw + 0.066
    return s


def slide_scale(prs):
    s = base_slide(prs, "ELASTIC SCALE",
                   "Absorbing a 10× spike without breaking a sweat")
    mech = [
        (TEAL, "Stateless hot path",
         ["Container Apps autoscale 3 → 60 replicas / region", "KEDA scales on HTTP concurrency in seconds"]),
        (BLUE, "Event-driven buffer",
         ["Event Hubs absorbs the surge; nothing dropped", "24h replay · Stream Analytics scales on depth"]),
        (PURPLE, "Data tier scales apart",
         ["Cosmos autoscale RU/s · Redis hot cache", "Reads stay < 2 ms under load"]),
        (GREEN, "Two regions, active/active",
         ["Front Door spreads load across SE + NEU", "A new region is a single flag"]),
    ]
    y = 2.0
    for accent, h, items in mech:
        column_card(s, 0.7, y, 5.7, 1.05, accent, h, items, body_size=11)
        y += 1.15
    # spike chart on the right
    cx, cy, cw, ch = 6.9, 2.0, 5.7, 3.55
    rect(s, cx, cy, cw, ch, fill=CARD, line=BORDER, lw=1.2, rounded=True)
    text(s, cx + 0.3, cy + 0.18, cw - 0.6, 0.4, "Throughput vs. latency during a spike",
         size=13, bold=True, color=TITLE)
    base_y = cy + ch - 0.7
    # bars (TPS) — baseline then spike then settle
    series = [5, 5, 5, 9, 16, 20, 14, 7, 5]
    bw = 0.42
    gap = 0.12
    bx = cx + 0.5
    maxv = 20
    for v in series:
        bh = (v / maxv) * 2.0
        rect(s, bx, base_y - bh, bw, bh, fill=(AMBER if v >= 16 else TEAL))
        bx += bw + gap
    # p99 latency flat line
    line_y = cy + 0.95
    rect(s, cx + 0.5, line_y, cw - 1.0, 0.03, fill=GREEN)
    text(s, cx + 0.5, line_y - 0.32, cw - 1.0, 0.3,
         "p99 latency stays flat (< 18 ms) — green line", size=10.5, color=GREEN)
    text(s, cx + 0.5, base_y + 0.06, cw - 1.0, 0.3,
         "5k baseline → 20k peak TPS (amber) absorbed in seconds", size=10.5, color=BODY)
    text(s, 0.7, 6.52, 11.9, 0.34,
         "Consumption pricing → capacity follows load and falls to near-zero off-peak.",
         size=12, color=HEAD, bold=True)
    return s


def slide_resilience(prs):
    s = base_slide(prs, "RESILIENCE & HIGH AVAILABILITY",
                   "What happens when something breaks")
    cols = [("Component", 0.2), ("If it fails", 0.34), ("What we do — blast radius", 0.46)]
    rows = [
        ["Scoring replica", "Pod crash / bad node",
         "ACA reschedules; peers serve; zero downtime (zone-redundant, min 3)"],
        ["Whole region", "Region outage",
         "Front Door fails over to North Europe; RTO < 5 min, RPO seconds"],
        ["Cosmos DB", "Primary unavailable",
         "Multi-master auto-failover; reads/writes continue in healthy region"],
        ["Redis cache", "Cache node down",
         "Cache-miss falls back to Cosmos point reads — degrade, don't fail"],
        ["Feature backend", "Cosmos/Redis unreachable",
         "Fail-safe: score on defaults + force step-up (SCA); never blind-approve, never 500"],
        ["Event Hubs", "Namespace impaired",
         "24h buffer + geo-DR alias; cold path catches up, no decision lost"],
        ["Azure OpenAI / agents", "LLM throttled / down",
         "Case work queues; the scoring hot path is unaffected (decoupled)"],
    ]
    accents = [TEAL, RED, BLUE, AMBER, GREEN, PURPLE, BLUE]
    table_block(s, 0.7, 1.95, 11.93, cols, rows, accents=accents, row_h=0.56)
    text(s, 0.7, 6.42, 11.93, 0.36,
         "Target: 99.99% availability · zone-redundant · dual-region active/active · every dependency has a fallback.",
         size=12, color=HEAD, bold=True)
    return s


def slide_security(prs):
    s = base_slide(prs, "SECURITY, ENFORCED",
                   "Not a policy on paper — enforced by the platform")
    cols = [("Control", 0.28), ("How it is enforced (mechanism)", 0.72)]
    rows = [
        ["No secrets in code",
         "Managed identities + Key Vault HSM; OIDC federation in CI; secret-scanning gate blocks the build"],
        ["EU regions only",
         "Azure Policy 'NordicSovereignty v3' denies any non-EU region at deploy time (deny effect)"],
        ["No public data planes",
         "Private endpoints on every PaaS; Policy denies public network access; WAF + mTLS at the edge"],
        ["Only Front Door reaches scoring",
         "X-Azure-FDID header validated in-app; ACA ingress locked; direct origin calls rejected"],
        ["Least privilege",
         "Entra ID + PIM just-in-time admin; scoped RBAC; workload identities; access reviews"],
        ["Provable posture",
         "Defender for Cloud CSPM/CWPP + Sentinel SIEM; STRIDE model 0 critical; full audit + Purview lineage"],
    ]
    accents = [TEAL, GREEN, BLUE, PURPLE, AMBER, GREEN]
    table_block(s, 0.7, 1.95, 11.93, cols, rows, accents=accents, row_h=0.64)
    text(s, 0.7, 6.48, 11.93, 0.34,
         "Every control has an owner, an enforcing mechanism, and a test in CI — security is verified, not assumed.",
         size=12, color=HEAD, bold=True)
    return s


def slide_ops_dashboard(prs):
    s = base_slide(prs, "OPERATIONS DASHBOARD", "What management sees — live")
    tiles = [
        (TEAL, "4,920 TPS", "Throughput", "12 / 60 replicas active"),
        (BLUE, "14 ms", "Scoring p99", "SLO < 18 ms · green"),
        (GREEN, "99.99%", "Availability 30d", "2 EU regions active/active"),
        (PURPLE, "€312k", "Fraud caught today", "1,204 cases opened"),
        (GREEN, "90 / 8 / 2", "Approve / SCA / Decline", "decision mix, %"),
        (AMBER, "1.1%", "False-positive rate", "▼ from 2.8%"),
        (BLUE, "18", "HITL review queue", "expert-review cases"),
        (TEAL, "auto", "EBA report", "next in 12 days"),
    ]
    x0, y0 = 0.7, 2.0
    cw, chh = 2.93, 1.6
    for i, (accent, v, l, d) in enumerate(tiles):
        col = i % 4
        row = i // 4
        x = x0 + col * (cw + 0.12)
        y = y0 + row * (chh + 0.18)
        rect(s, x, y, cw, chh, fill=CARD, line=BORDER, lw=1.0, rounded=True)
        rect(s, x + 0.12, y, cw - 0.24, 0.05, fill=accent)
        text(s, x + 0.05, y + 0.2, cw - 0.1, 0.6, v, size=23, bold=True,
             color=accent, align=CENTER)
        text(s, x + 0.1, y + 0.92, cw - 0.2, 0.35, l, size=12, bold=True,
             color=TITLE, align=CENTER)
        text(s, x + 0.1, y + 1.24, cw - 0.2, 0.3, d, size=9.5, color=BODY,
             align=CENTER)
    text(s, 0.7, 5.78, 11.93, 0.85, [
        ("Built into the demo console (open /ops) and mirrored in Power BI — the same numbers, exec-friendly.",
         {"color": HEAD, "bold": True, "size": 13, "space": 4}),
        ("Auto-refreshes every 2 seconds · drill-through from any KPI to a single decision and its full trace.",
         {"color": BODY, "size": 12}),
    ])
    return s


def slide_roadmap(prs):
    s = base_slide(prs, "ROADMAP & MILESTONES",
                   "A 12-month path, with milestones that prove value")
    # timeline spine
    ty = 3.85
    rect(s, 0.9, ty, 11.5, 0.04, fill=BORDER)
    miles = [
        (TEAL, "M1 · Wk 2", "Pilot live", "Sweden Central · shadow scoring on live traffic", True),
        (BLUE, "M2 · Wk 6", "Value proven", "−41% validated on your data · security sign-off", False),
        (PURPLE, "M3 · Wk 12", "Production", "Active/active DR · EBA reporting automated", True),
        (AMBER, "M4 · Mo 5", "Optimised", "SCA optimiser tuned · exemptions 22% → 73%", False),
        (GREEN, "M5 · Mo 9", "Five markets", "SE·NO·DK·FI·EE live · agentic automation expanded", True),
        (HEAD, "M6 · Mo 12", "Assured", "Continuous governance · ISO 42001 track", False),
    ]
    n = len(miles)
    x0 = 1.1
    span = 11.1
    for i, (accent, when, head, desc, above) in enumerate(miles):
        cx = x0 + span * (i / (n - 1))
        # node
        dot = rect(s, cx - 0.09, ty - 0.07, 0.18, 0.18, fill=accent, rounded=True)
        if above:
            cardy = ty - 1.55
        else:
            cardy = ty + 0.32
        bx = min(max(cx - 0.95, 0.7), 11.73)
        rect(s, bx, cardy, 1.9, 1.2, fill=CARD, line=BORDER, lw=1.0, rounded=True)
        rect(s, bx + 0.1, cardy, 1.7, 0.05, fill=accent)
        text(s, bx + 0.1, cardy + 0.12, 1.7, 0.3, when, size=10.5, bold=True, color=accent)
        text(s, bx + 0.1, cardy + 0.4, 1.7, 0.3, head, size=12, bold=True, color=TITLE)
        text(s, bx + 0.1, cardy + 0.68, 1.7, 0.5, desc, size=8.8, color=BODY, space=1)
    text(s, 0.7, 6.5, 11.9, 0.32,
         "Each milestone has an exit criterion: we only advance when the prior value is proven on your data.",
         size=12, color=HEAD, bold=True)
    return s


# --------------------------------------------------------------------------- #
# Personalisation footer
# --------------------------------------------------------------------------- #
def add_footer(s, page_no, dark=True):
    fg = BODY
    # monogram chip
    rect(s, 0.7, 7.04, 0.46, 0.34, fill=TEAL, rounded=True)
    text(s, 0.7, 7.06, 0.46, 0.3, "JRP", size=11, bold=True, color=NAVY,
         align=CENTER, anchor=MID)
    text(s, 1.24, 7.04, 5.0, 0.34, f"{AUTHOR} · {ROLE}", size=9, color=fg, anchor=MID)
    text(s, 7.0, 7.04, 4.8, 0.34, BRAND, size=9, color=fg, align=PP_ALIGN.RIGHT, anchor=MID)
    text(s, 11.95, 7.04, 0.68, 0.34, str(page_no), size=9, bold=True, color=TEAL,
         align=PP_ALIGN.RIGHT, anchor=MID)


# --------------------------------------------------------------------------- #
# Assemble
# --------------------------------------------------------------------------- #
def reorder(prs, order):
    sldIdLst = prs.slides._sldIdLst
    elems = list(sldIdLst)
    for e in elems:
        sldIdLst.remove(e)
    for i in order:
        sldIdLst.append(elems[i])


def main():
    prs = Presentation(SRC)
    n_orig = len(prs.slides)  # 34

    # Hero / full-art slides that should NOT get a footer (by original index).
    hero_idx = {0, 1, 20, 26, 27, 28}
    hero_ids = {prs.slides[i].slide_id for i in hero_idx}

    # Build new slides (appended at the end, indices n_orig..n_orig+7).
    A = slide_inputs_outputs(prs)   # after orig 8
    B = slide_where_fits(prs)       # after orig 10
    C = slide_threats(prs)          # after orig 17
    D = slide_scale(prs)            # after orig 17 (follows C)
    F = slide_security(prs)         # after orig 23
    G = slide_ops_dashboard(prs)    # after orig 24
    E = slide_resilience(prs)       # after orig 24 (follows G)
    H = slide_roadmap(prs)          # after orig 25

    idx = {id(x): n_orig + k for k, x in enumerate([A, B, C, D, F, G, E, H])}
    inserts = {
        8: [A],
        10: [B],
        17: [C, D],
        23: [F],
        24: [G, E],
        25: [H],
    }
    order = []
    for i in range(n_orig):
        order.append(i)
        for ns in inserts.get(i, []):
            order.append(idx[id(ns)])
    reorder(prs, order)

    # Footer + page numbers on the final ordering. Skip a slide if its existing
    # content already reaches into the footer band (avoid any overlap).
    from pptx.util import Emu as _Emu

    def lowest_text_bottom(slide):
        low = 0.0
        for sh in slide.shapes:
            if sh.top is None or not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            b = _Emu(sh.top).inches + _Emu(sh.height).inches
            low = max(low, b)
        return low

    for pos, s in enumerate(prs.slides, start=1):
        if s.slide_id in hero_ids:
            continue
        if lowest_text_bottom(s) > 6.95:
            continue
        add_footer(s, pos)

    prs.save(OUT)
    print(f"WROTE {OUT}")
    print(f"slides: {len(prs.slides)} (was {n_orig}, added {len(prs.slides) - n_orig})")


if __name__ == "__main__":
    main()
