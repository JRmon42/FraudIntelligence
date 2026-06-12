"""
Build the 45-min AMA Capstone deck for Case Study 30 (Heimdall)
using the supplied Azure Master Architect template.

Structure (target ~22 slides for a 45-min talk with demo + Q&A):
 1  Title
 2  Agenda
 3  The customer & the problem (3 columns: who/what/why)
 4  Business case (3 columns: cost of fraud / cost of false declines / regulatory pressure)
 5  Target outcomes (3 columns KPIs)
 6  Solution at a glance (right content + image placeholder)
 7  High-level architecture (image)
 8  Hot path: real-time scoring <18ms (right content + image)
 9  Cold path: graph + analytics (right content)
10  AI/ML strategy: ensemble + GNN + LLM (3 columns)
11  Agentic AI orchestration (right content)
12  Compliance: GDPR, EU AI Act, PSD2, EBA (3 columns)
13  Sovereignty for the Nordics (right content)
14  Security & DevSecOps (right content)
15  MLOps & responsible AI (right content)
16  Multi-region resilience (right content)
17  FinOps: cost-to-serve & TCO (right content)
18  Live demo (Title only)
19  Key benefits & value (Content w/image_2)
20  Risks & mitigations (3 columns)
21  Roadmap / next steps (3 columns)
22  Closing
"""
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
from pptx.oxml.ns import qn

SRC = "/tmp/slidebuild/template.pptx"
OUT = "/mnt/c/Users/jpontvianne/Documents/Azure/MasterArchitect/Code/Heimdall/slides/Heimdall_AMA_Capstone.pptx"

p = Presentation(SRC)

# Wipe template-provided slides; we'll rebuild from layouts
sldIdLst = p.slides._sldIdLst
for sld in list(sldIdLst):
    rId = sld.get(qn("r:id"))
    p.part.drop_rel(rId)
    sldIdLst.remove(sld)

LAYOUTS = {l.name: l for l in p.slide_layouts}

def set_text(ph, text, font_size=None):
    tf = ph.text_frame
    tf.clear()
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        if font_size:
            for r in para.runs:
                r.font.size = Pt(font_size)

def add(layout_name, title=None, body_map=None, notes=None):
    """body_map: dict {placeholder_idx: text_or_list}"""
    layout = LAYOUTS[layout_name]
    s = p.slides.add_slide(layout)
    if title is not None:
        # title placeholder is idx=0
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 0:
                set_text(ph, title)
                break
    if body_map:
        for idx, txt in body_map.items():
            for ph in s.placeholders:
                if ph.placeholder_format.idx == idx:
                    set_text(ph, txt)
                    break
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s

# --- 1 Title ---
add("Title slide", title="Heimdall — AI-Driven Heimdall Platform\nNordic Payments Provider · AMA Capstone — Case Study 30",
    notes=("Welcome. Over the next 45 minutes I will walk you through how we deliver real-time, sovereign, "
           "agentic fraud intelligence on Azure for a Stockholm-based payments provider processing 4.2 billion "
           "transactions a year across the Nordics and Baltics. The session covers business case, architecture, "
           "AI strategy, compliance, a live demo, and roadmap. I'm happy to take questions throughout."))

# --- 2 Agenda ---
add("Agenda", title="AGENDA",
    body_map={
        10: "Customer context & the business case",
        11: "Target outcomes",
        12: "Reference architecture",
        13: "Real-time scoring (<18 ms)",
        14: "AI / ML & agentic orchestration",
        15: "Compliance, sovereignty & security",
        16: "Live demo",
        17: "Roadmap, value, Q&A",
    },
    notes="I will keep each section to roughly 5 minutes, with 10 minutes for the demo and 5 for Q&A.")

# --- 3 Customer & problem ---
add("3 Columns", title="The customer & the problem",
    body_map={
        10: ("Stockholm-based payments infrastructure provider. 4.2 B transactions per year across "
             "Sweden, Norway, Denmark, Finland, Estonia. Operates payment rails for hundreds of issuers, "
             "acquirers and merchants in the Nordic-Baltic corridor."),
        11: ("Fraud losses are growing 34 % year-on-year and outpace detection capability. False declines run at 2.8 %, "
             "eroding merchant trust and driving cardholder churn. Real-time decisioning must complete inside 120 ms "
             "but the legacy stack cannot meet it."),
        12: ("PSD2 SCA exemption management is manual, EBA fraud reporting takes 320 person-hours per quarter, "
             "and the EU AI Act will classify fraud scoring as a high-risk system from August 2026."),
        26: "WHO",
        27: "WHAT",
        28: "WHY NOW",
    },
    notes="Anchor the audience: this is a regulated, latency-critical, multi-country problem — three things our reference architecture must solve simultaneously.")

# --- 4 Business case ---
add("3 Columns", title="The business case",
    body_map={
        10: ("Fraud losses today: ~€48 M / year. A 41 % reduction returns ~€19.7 M annually with a payback period of "
             "under 7 months on the platform investment."),
        11: ("False decline cost: ~€31 M / year in lost revenue and CLV erosion. Cutting 2.8 % → 1.1 % returns "
             "~€19 M and protects merchant retention."),
        12: ("Regulatory cost & risk: EBA fines for under-reporting average €4 M; EU AI Act non-conformity penalties up to "
             "€35 M or 7 % of group turnover. Automation removes the exposure and frees the compliance team."),
        26: "FRAUD",
        27: "FRICTION",
        28: "REGULATION",
    },
    notes="Lead with money to engage the CEO and CFO; risk to engage the CISO; throughput/SLA to engage the CTO.")

# --- 5 Target outcomes ---
add("3 Columns", title="Target outcomes",
    body_map={
        10: ("Fraud loss reduced by 41 %.\nDecline rate reduced from 2.8 % to 1.1 %.\nReal-time scoring p99 < 18 ms (target), "
             "measured 14 ms in load test at 5 k TPS sustained, 20 k TPS peak."),
        11: ("PSD2 SCA exemption coverage lifted from 22 % to 73 % of eligible transactions while remaining inside the "
             "EBA fraud-rate corridor.\nEBA fraud reporting fully automated — quarterly cycle from 320 hours to zero."),
        12: ("EU AI Act high-risk system conformity from day one: model cards, data governance, human oversight, "
             "logging, post-market monitoring all built in."),
        26: "PERFORMANCE",
        27: "COMPLIANCE",
        28: "RESPONSIBLE AI",
    })

# --- 6 Solution at a glance ---
add("Right Content", title="Solution at a glance",
    body_map={
        10: ("A real-time, multi-region, sovereign AI platform on Azure.\n\n"
             "• Hot path: Front Door → Container Apps (FastAPI + ONNX) → Cosmos DB · target p99 < 18 ms\n"
             "• Cold path: Event Hubs → Stream Analytics → Cosmos Gremlin → Microsoft Fabric medallion → Power BI\n"
             "• AI: stacked ensemble + Graph Neural Network + Azure OpenAI for narrative / case reasoning\n"
             "• Agentic layer: Microsoft Semantic Kernel orchestrating six specialised agents with reflection\n"
             "• Governance: Microsoft Purview, Defender for Cloud, Azure Policy, customer-managed keys"),
    })

# --- 7 Architecture (image placeholder, narrative on slide) ---
add("Title Only", title="Reference architecture",
    notes=("The architecture is split into a hot synchronous path optimised for sub-18 ms decisioning and a cold asynchronous "
           "path that powers learning, reporting and the agentic case workflow. Front Door fronts both, with WAF and bot "
           "protection. The scoring API is single-purpose, stateless, and runs ONNX Runtime in-process. Decisions are "
           "fire-and-forget published to Event Hubs, Stream Analytics computes rolling features, Cosmos DB serves both as "
           "feature store and Gremlin graph for ring detection, and Fabric powers the medallion lakehouse that feeds Power "
           "BI and the EBA reporter."))

# --- 8 Real-time scoring ---
add("Right Content", title="Hot path — real-time scoring under 18 ms",
    body_map={
        10: ("FastAPI · Python 3.11 · ONNX Runtime in-process · Azure Container Apps Dedicated D8 workload profile.\n\n"
             "• Cosmos DB point reads (<2 ms p99) feed card & merchant features\n"
             "• Async LRU hot cache for top 1 M cards (60 s TTL)\n"
             "• Stacked ensemble (XGBoost + LightGBM + Logistic) calibrated with isotonic regression\n"
             "• PSD2 exemption optimiser selects the best applicable exemption per transaction\n"
             "• Decision and explanation emitted to Event Hubs (decision.events) for downstream learning\n"
             "• OTEL spans per stage; load-tested at 5 k TPS sustained with p99 = 14 ms")
    })

# --- 9 Cold path ---
add("Right Content", title="Cold path — graph & analytics",
    body_map={
        10: ("Event Hubs (geo-DR) is the single source of truth for transaction telemetry.\n\n"
             "• Stream Analytics computes 1 m / 5 m / 1 h / 24 h rolling aggregates and writes to Cosmos features\n"
             "• Cosmos Gremlin holds the heterogeneous card / merchant / device graph used by the GNN\n"
             "• Microsoft Fabric (OneLake) lakehouse — Bronze ingest from EH Capture, Silver clean & tokenise, Gold marts\n"
             "• EBA quarterly fraud report generated automatically by a containerised job feeding Power BI Premium")
    })

# --- 10 AI/ML ---
add("3 Columns", title="AI / ML strategy",
    body_map={
        10: ("Stacked ensemble: XGBoost + LightGBM + Logistic Regression with a meta-learner. Calibrated, exported to "
             "ONNX for sub-5 ms in-process scoring. Re-trained nightly via Azure ML pipelines."),
        11: ("Graph Neural Network: PyTorch Geometric, GraphSAGE on a heterogeneous card / merchant / device / IP graph. "
             "Outputs node embeddings + ring-membership probability. Surfaces synthetic identity & merchant-collusion patterns."),
        12: ("Azure OpenAI: gpt-4o-mini for low-latency reasoning and tool calling, gpt-4o for narrative SAR / EBA drafting. "
             "Prompt-versioned, evaluated for groundedness and toxicity, monitored for drift."),
        26: "ENSEMBLE", 27: "GRAPH", 28: "GENERATIVE",
    })

# --- 11 Agentic AI ---
add("Right Content", title="Agentic AI — multi-agent case workflow",
    body_map={
        10: ("Built on Microsoft Semantic Kernel with a state-graph planner and a reflector loop.\n\n"
             "• TriageAgent classifies the alert and routes\n"
             "• GraphAnalystAgent traverses Cosmos Gremlin two-hop neighbourhoods\n"
             "• PolicyAgent maps findings to PSD2 / EBA rules via tool calling\n"
             "• CaseManagerAgent persists timeline & decisions in Cosmos\n"
             "• NarrativeAgent drafts the SAR / EBA narrative\n"
             "• ReflectorAgent reviews each step and can re-route — true autonomy with bounded reflection budget\n\n"
             "Demonstrates planning, handoffs, tool use, reflection, and human-in-the-loop checkpoints.")
    })

# --- 12 Compliance ---
add("3 Columns", title="Compliance by design",
    body_map={
        10: ("GDPR — Article 22 automated decision controls, Article 25 privacy by design, Article 32 security, Article 35 DPIA "
             "completed, Article 44 transfer minimisation."),
        11: ("EU AI Act — high-risk system. Risk management (Art 9), data governance (Art 10), technical documentation (Art 11), "
             "logging (Art 12), transparency (Art 13), human oversight (Art 14), accuracy & cybersecurity (Art 15)."),
        12: ("PSD2 SCA exemption framework + EBA Guidelines on fraud reporting (EBA/GL/2020/01) — fully automated quarterly "
             "report by instrument type."),
        26: "DATA", 27: "AI ACT", 28: "PAYMENTS",
    })

# --- 13 Sovereignty ---
add("Right Content", title="Sovereignty for SE / NO / DK / FI / EE",
    body_map={
        10: ("• Primary region Sweden Central; DR North Europe — both EU jurisdictions\n"
             "• Azure Policy enforces allowed locations and denies any non-EU placement\n"
             "• Customer-managed keys in Key Vault (HSM-backed); BYOK rotation every 90 days\n"
             "• Microsoft Purview classifies and labels PII; tokenisation in Silver layer of the lakehouse\n"
             "• Confidential VMs available for AML compute on demand\n"
             "• Audit log immutability via Storage immutable blob policies; logs to a sovereign Log Analytics workspace\n"
             "• Country-specific reporting filters in Power BI honour local supervisor requirements")
    })

# --- 14 Security & DevSecOps ---
add("Right Content", title="Security & DevSecOps",
    body_map={
        10: ("• Defender for Cloud (Servers P2, Containers, KeyVault, Storage, CosmosDB, AppServices, OpenAI)\n"
             "• Front Door Premium WAF, OWASP 3.2, bot manager, mTLS to scoring API\n"
             "• Private endpoints on every PaaS service, public network access disabled\n"
             "• Managed identities everywhere — zero secrets in code; secrets exclusively in Key Vault\n"
             "• OIDC federated GitHub Actions deploy; Bicep validated, CodeQL & Trivy in CI; signed container images (cosign)\n"
             "• SBOM produced per release; Dependabot + Defender for DevOps")
    })

# --- 15 MLOps & responsible AI ---
add("Right Content", title="MLOps & responsible AI",
    body_map={
        10: ("• Azure ML workspaces (prod / dev) with managed network, MLflow tracking, model registry\n"
             "• Pipelines: data → train → evaluate → register → online endpoint with traffic split (canary 5 %)\n"
             "• Drift detection (data + concept) on Cosmos features; auto-trigger retrain when drift > threshold\n"
             "• Model cards + intended use + limitations + fairness audit per model (EU AI Act Art 11 evidence)\n"
             "• Per-country fairness monitoring — disparity in TPR < 5 % is the gate\n"
             "• Human review queue for high-impact decisions; agentic workflow logs every reasoning step")
    })

# --- 16 Multi-region resilience ---
add("Right Content", title="Multi-region resilience",
    body_map={
        10: ("• Cosmos DB multi-region writes, session consistency, automatic regional failover\n"
             "• Event Hubs geo-DR pairing with namespace alias — RPO seconds, RTO < 5 min\n"
             "• ACR Premium geo-replicated; Container Apps deployed to both regions, fronted by AFD with weighted routing\n"
             "• AML workspaces per region; model registry replicated\n"
             "• Tested DR runbook in `docs/runbook.md`, quarterly fire-drill schedule")
    })

# --- 17 FinOps ---
add("Right Content", title="FinOps — cost to serve",
    body_map={
        10: ("• Cost-to-serve target: < €0.0008 per scored transaction at 4.2 B / year (~€3.4 M / year all-in)\n"
             "• Reserved instances on AML CPU pools, Cosmos autoscale capped, Fabric F2 capacity, AFD Standard tier where allowed\n"
             "• Cost guard: scripts/scale-to-min.sh pauses Fabric, scales AML to zero, halts ASA between demos\n"
             "• Tagging convention enforced by Azure Policy; weekly Cost Management dashboard, alerts at 80 / 100 % of budget")
    })

# --- 18 Demo ---
add("Title Only", title="Live demo",
    notes=("Demo plan (10 minutes):\n"
           "1. Show Power BI executive dashboard at idle.\n"
           "2. Run transaction-simulator at 2 k TPS with --pattern mixed; show p99 stays <18 ms in Grafana.\n"
           "3. Inject --pattern fraud-ring (10 cards / 3 merchants / circular flow); GNN catches the ring; alert fires.\n"
           "4. Open the agentic console; watch TriageAgent → GraphAnalystAgent → PolicyAgent → CaseManagerAgent → NarrativeAgent.\n"
           "5. Open the auto-generated EBA Q1 report in Power BI."))

# --- 19 Key benefits & value (image placeholder layout) ---
add("Content w/image_2", title="Key benefits & value delivery",
    body_map={
        10: ("• €38.7 M annualised value (fraud + false-decline reduction)\n"
             "• 320 → 0 person-hours per quarter on EBA reporting\n"
             "• 73 % SCA exemption coverage — measurable cardholder UX uplift\n"
             "• EU AI Act high-risk conformity from day one — no remediation backlog\n"
             "• Platform extensible to AML, dispute management and KYC re-verification — same agentic scaffold\n"
             "• Regulator-ready audit trail by construction — full lineage from raw event to decision and narrative")
    })

# --- 20 Risks & mitigations ---
add("3 Columns", title="Risks & mitigations",
    body_map={
        10: ("Latency regression as feature volume grows.\nMitigation: continuous load tests in CI; budget alerts on p99; "
             "ONNX graph optimisation; ACA dedicated workload profiles; hot card cache."),
        11: ("Model bias across the five Nordic markets.\nMitigation: per-country fairness gate, monthly audit, drift "
             "detection auto-retrain, human review queue for high-impact decisions."),
        12: ("Vendor / region availability.\nMitigation: multi-region active-active, geo-DR, runbook tested quarterly, "
             "fallback rules-based scorer for graceful degradation."),
        26: "PERFORMANCE", 27: "FAIRNESS", 28: "RESILIENCE",
    })

# --- 21 Roadmap ---
add("3 Columns", title="Roadmap & next steps",
    body_map={
        10: ("Now → 90 days: complete pilot in Sweden Central, regulator engagement on EU AI Act conformity package, "
             "blue/green of the legacy scorer."),
        11: ("90 → 180 days: roll-out across Norway, Denmark, Finland, Estonia; activate DR; onboard first issuer customer "
             "to the agentic case console; integrate KYC re-verification agent."),
        12: ("180 days+: extend platform to AML transaction monitoring & dispute orchestration; formalise model marketplace; "
             "publish responsible-AI report; pursue ISO 42001 certification."),
        26: "STABILISE", 27: "SCALE", 28: "EXTEND",
    })

# --- 22 Closing ---
add("1_Closing logo slide",
    notes="Thank you. Questions, challenges and counter-arguments very welcome — that is the point of this cohort.")

p.save(OUT)
print("WROTE", OUT, "slides:", len(p.slides))
